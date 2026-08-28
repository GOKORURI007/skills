#!/usr/bin/env python3
"""Focused regressions for chart geometry and image fitting."""

import importlib.util
import struct
import sys
import tempfile
import unittest
import zlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

SCRIPT = Path(__file__).resolve().parents[1] / "native_svg_to_ppt.py"
SPEC = importlib.util.spec_from_file_location("native_svg_to_ppt", SCRIPT)
CONVERTER = importlib.util.module_from_spec(SPEC)
sys.modules[SPEC.name] = CONVERTER
SPEC.loader.exec_module(CONVERTER)


def png_bytes(width, height):
    raw = b"".join(b"\x00" + b"\x33\x77\xaa" * width for _ in range(height))

    def chunk(kind, payload):
        return struct.pack(">I", len(payload)) + kind + payload + struct.pack(">I", zlib.crc32(kind + payload) & 0xFFFFFFFF)

    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )


class ConverterGeometryTests(unittest.TestCase):
    def test_svg_arc_uses_real_ellipse_geometry(self):
        points = CONVERTER._parse_path_to_points("M 100 100 A 50 50 0 0 1 200 100")
        arc = points[0]
        self.assertGreater(len(arc), 20)
        self.assertAlmostEqual(arc[-1][0], 200.0, places=5)
        self.assertAlmostEqual(arc[-1][1], 100.0, places=5)
        self.assertGreater(max(abs(y - 100.0) for _, y in arc), 45.0)

    def test_transform_order_and_nested_translation(self):
        self.assertEqual(CONVERTER._parse_axis_aligned_transform("translate(100 50) scale(2)"), (100.0, 50.0, 2.0, 2.0))
        self.assertEqual(CONVERTER._parse_axis_aligned_transform("scale(2) translate(100 50)"), (200.0, 100.0, 2.0, 2.0))

    def test_picture_meet_and_slice_preserve_ratio(self):
        with tempfile.TemporaryDirectory() as temp:
            image_path = Path(temp) / "wide.png"
            image_path.write_bytes(png_bytes(200, 100))
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            meet = CONVERTER.add_fitted_picture(slide, str(image_path), 0, 0, 2, 2, "xMidYMid meet")
            self.assertAlmostEqual(meet.width / meet.height, 2.0, places=2)
            sliced = CONVERTER.add_fitted_picture(slide, str(image_path), 3, 0, 2, 2, "xMidYMid slice")
            self.assertEqual(sliced.width, Inches(2))
            self.assertEqual(sliced.height, Inches(2))
            self.assertGreater(sliced.crop_left, 0)
            self.assertGreater(sliced.crop_right, 0)


if __name__ == "__main__":
    unittest.main()
