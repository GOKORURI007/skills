#!/usr/bin/env python3
"""
PPTX Template Asset Extractor for Planner's PPT Hell v2.

Adds structural candidates to template_profile.json with:
  - Color values from ALL slide masters (XML-parsed)
  - Font families from ALL slide masters (XML-parsed)
  - Bitmap assets from ALL slides (SHA256-deduped)
  - Font size measurements from svg-flat/ SVGs (when available)
  - Decoration patterns from svg-flat/ SVGs (when available)
These candidates never replace the visual Template Worker's conclusions or
become binding without its page-by-page confirmation.

Usage:
  python3 extract_template_assets.py <source.pptx> --project <project_dir>

Requires: python-pptx >= 0.6.21
"""

import argparse
import hashlib
import json
import os
import re
import shutil
import sys
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

EMU_PER_INCH = 914400
DEFAULT_DPI = 96


def emu_to_px(emu, dpi=DEFAULT_DPI):
    """Convert EMU (English Metric Unit) to pixels at given DPI."""
    if emu is None:
        return 0.0
    return round(emu / EMU_PER_INCH * dpi, 1)


def sha256_blob(blob):
    """Return SHA-256 hex digest of a bytes blob."""
    return hashlib.sha256(blob).hexdigest()


# ---------------------------------------------------------------------------
# a) extract_theme_colors_all_masters
# ---------------------------------------------------------------------------

_COLOR_ROLE_MAP = {
    "accent1": "primary",
    "accent2": "accent",
    "accent3": "accent-tertiary",
    "accent4": "accent-quaternary",
    "accent5": "accent-quinary",
    "accent6": "accent-senary",
    "dk1": "text-primary",
    "dk2": "text-secondary",
    "lt1": "background",
    "lt2": "surface",
    "hlink": "link",
    "folHlink": "link-visited",
}


def _read_theme_xmls_from_zip(pptx_path):
    """
    Read all theme XML files from the PPTX zip archive.
    Returns list of ElementTree objects for each theme*.xml found.
    """
    import zipfile
    import xml.etree.ElementTree as ET
    themes = []
    try:
        with zipfile.ZipFile(str(pptx_path), 'r') as z:
            theme_files = sorted([f for f in z.namelist() if f.startswith('ppt/theme/') and f.endswith('.xml')])
            for tf in theme_files:
                try:
                    content = z.read(tf)
                    root = ET.fromstring(content)
                    themes.append(root)
                except Exception:
                    continue
    except Exception:
        pass
    return themes


def _parse_clr_scheme_from_theme(theme_root):
    """
    Parse a single theme XML root for clrScheme.
    Returns dict mapping role -> hex.
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    result = {}
    try:
        clr_scheme = theme_root.find(f'.//{{{ns}}}clrScheme')
        if clr_scheme is None:
            return result
    except Exception:
        return result

    for name, role in _COLOR_ROLE_MAP.items():
        try:
            elem = clr_scheme.find(f'{{{ns}}}{name}')
            if elem is None:
                continue
            srgb = elem.find(f'{{{ns}}}srgbClr')
            if srgb is not None:
                hex_val = '#' + srgb.get('val', '000000')
            else:
                sys_clr = elem.find(f'{{{ns}}}sysClr')
                if sys_clr is not None:
                    hex_val = '#' + sys_clr.get('lastClr', '000000')
                else:
                    continue
            result[role] = hex_val
        except Exception:
            continue
    return result


def _parse_font_scheme_from_theme(theme_root):
    """
    Parse a single theme XML root for fontScheme.
    Returns dict with heading and body font_family.
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    result = {}
    try:
        font_scheme = theme_root.find(f'.//{{{ns}}}fontScheme')
        if font_scheme is None:
            return result
    except Exception:
        return result

    try:
        major = font_scheme.find(f'{{{ns}}}majorFont')
        if major is not None:
            ea = major.find(f'{{{ns}}}ea')
            latin = major.find(f'{{{ns}}}latin')
            heading_font = ""
            if ea is not None and ea.get('typeface'):
                heading_font = ea.get('typeface')
            elif latin is not None and latin.get('typeface'):
                heading_font = latin.get('typeface')
            if heading_font:
                result["heading"] = heading_font
    except Exception:
        pass

    try:
        minor = font_scheme.find(f'{{{ns}}}minorFont')
        if minor is not None:
            ea = minor.find(f'{{{ns}}}ea')
            latin = minor.find(f'{{{ns}}}latin')
            body_font = ""
            if ea is not None and ea.get('typeface'):
                body_font = ea.get('typeface')
            elif latin is not None and latin.get('typeface'):
                body_font = latin.get('typeface')
            if body_font:
                result["body"] = body_font
    except Exception:
        pass

    return result


def extract_theme_colors_all_masters(prs, pptx_path=None):
    """
    Iterate ALL slide masters and read clrScheme XML.
    Falls back to reading theme XMLs directly from the PPTX zip if masters
    don't inline the theme.

    Returns list of {role, hex, source:"xml_parsed", confidence:"high", note:""}.
    Takes the first master's theme when multiple themes define the same role.
    """
    colors = []
    seen_roles = set()

    # First try: via slide master element (works when theme is inlined)
    for master in prs.slide_masters:
        try:
            theme_el = master.element.find('.//' + qn('a:clrScheme'))
            if theme_el is None:
                continue
        except Exception:
            continue

        for name, role in _COLOR_ROLE_MAP.items():
            if role in seen_roles:
                continue
            try:
                elem = theme_el.find(qn(f'a:{name}'))
                if elem is None:
                    continue
                srgb = elem.find(qn('a:srgbClr'))
                if srgb is not None:
                    hex_val = '#' + srgb.get('val', '000000')
                else:
                    sys_clr = elem.find(qn('a:sysClr'))
                    if sys_clr is not None:
                        hex_val = '#' + sys_clr.get('lastClr', '000000')
                    else:
                        continue

                colors.append({
                    "role": role,
                    "hex": hex_val,
                    "source": "xml_parsed",
                    "confidence": "high",
                    "note": "",
                })
                seen_roles.add(role)
            except Exception:
                continue

    # Second try: read theme XMLs directly from PPTX zip
    if not colors and pptx_path:
        themes = _read_theme_xmls_from_zip(pptx_path)
        for theme_root in themes:
            scheme = _parse_clr_scheme_from_theme(theme_root)
            for role, hex_val in scheme.items():
                if role not in seen_roles:
                    colors.append({
                        "role": role,
                        "hex": hex_val,
                        "source": "xml_parsed",
                        "confidence": "high",
                        "note": "",
                    })
                    seen_roles.add(role)

    return colors


# ---------------------------------------------------------------------------
# b) extract_theme_fonts_all_masters
# ---------------------------------------------------------------------------

def extract_theme_fonts_all_masters(prs, pptx_path=None):
    """
    Iterate ALL slide masters and read fontScheme XML.
    Falls back to reading theme XMLs directly from the PPTX zip if masters
    don't inline the theme.

    majorFont -> heading role, minorFont -> body role.
    Prefer EA (East Asian) font for CJK contexts.
    Returns list of {role, font_family, source:"xml_parsed", confidence:"high"}.
    """
    fonts = []
    seen_roles = set()

    # First try: via slide master element
    for master in prs.slide_masters:
        try:
            font_scheme = master.element.find('.//' + qn('a:fontScheme'))
            if font_scheme is None:
                continue
        except Exception:
            continue

        # majorFont -> heading
        if "heading" not in seen_roles:
            try:
                major = font_scheme.find(qn('a:majorFont'))
                if major is not None:
                    ea = major.find(qn('a:ea'))
                    latin = major.find(qn('a:latin'))
                    font_family = ""
                    if ea is not None and ea.get('typeface'):
                        font_family = ea.get('typeface')
                    elif latin is not None and latin.get('typeface'):
                        font_family = latin.get('typeface')
                    if font_family:
                        fonts.append({
                            "role": "heading",
                            "font_family": font_family,
                            "source": "xml_parsed",
                            "confidence": "high",
                        })
                        seen_roles.add("heading")
            except Exception:
                pass

        # minorFont -> body
        if "body" not in seen_roles:
            try:
                minor = font_scheme.find(qn('a:minorFont'))
                if minor is not None:
                    ea = minor.find(qn('a:ea'))
                    latin = minor.find(qn('a:latin'))
                    font_family = ""
                    if ea is not None and ea.get('typeface'):
                        font_family = ea.get('typeface')
                    elif latin is not None and latin.get('typeface'):
                        font_family = latin.get('typeface')
                    if font_family:
                        fonts.append({
                            "role": "body",
                            "font_family": font_family,
                            "source": "xml_parsed",
                            "confidence": "high",
                        })
                        seen_roles.add("body")
            except Exception:
                pass

    # Second try: read theme XMLs directly from PPTX zip
    if not fonts and pptx_path:
        themes = _read_theme_xmls_from_zip(pptx_path)
        for theme_root in themes:
            scheme = _parse_font_scheme_from_theme(theme_root)
            for role, font_family in scheme.items():
                if role not in seen_roles:
                    fonts.append({
                        "role": role,
                        "font_family": font_family,
                        "source": "xml_parsed",
                        "confidence": "high",
                    })
                    seen_roles.add(role)

    return fonts


# ---------------------------------------------------------------------------
# c) extract_assets_from_slides
# ---------------------------------------------------------------------------

def _extract_image_blob(prs, rId):
    """Extract image bytes from a relationship ID in the PPTX package."""
    try:
        part = prs.part.related_part(rId)
        return part.blob
    except Exception:
        return None


def _image_ext_from_content_type(content_type):
    """Map content type to file extension."""
    mapping = {
        "image/jpeg": ".jpeg",
        "image/png": ".png",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/webp": ".webp",
        "image/svg+xml": ".svg",
    }
    return mapping.get(content_type, ".bin")


def extract_assets_from_slides(prs, project_dir):
    """
    Iterate ALL slides and extract picture shapes (shape_type == 13 / Picture).
    Copy images to template_media/<asset_id>.<ext> with SHA256 dedup.
    Detect full_bleed_background when shape covers nearly the entire slide.
    Detect reused images (same SHA256 blob on multiple slides).

    Returns list of asset dicts.
    """
    media_dir = Path(project_dir) / "_internal" / "00_project" / "template_media"
    media_dir.mkdir(parents=True, exist_ok=True)

    canvas_w_px = emu_to_px(prs.slide_width)
    canvas_h_px = emu_to_px(prs.slide_height)
    canvas_area_px = canvas_w_px * canvas_h_px
    canvas_str = f"{int(canvas_w_px)}x{int(canvas_h_px)}"

    # Track assets: sha256 -> {asset_id, file, pages[], geometries[], blob}
    sha_assets = {}
    # Track asset_id counter
    asset_counters = {"bg": 0, "logo": 0, "img": 0, "deco": 0}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type != 13:  # Picture
                continue

            try:
                img = shape.image
                blob = img.blob
                if not blob:
                    continue
            except Exception:
                continue

            blob_sha = sha256_blob(blob)

            # Geometry
            sx = emu_to_px(shape.left)
            sy = emu_to_px(shape.top)
            sw = emu_to_px(shape.width)
            sh = emu_to_px(shape.height)

            # Determine content type and extension
            content_type = getattr(img, 'content_type', 'image/jpeg')
            ext = _image_ext_from_content_type(content_type)

            area_ratio = (sw * sh) / canvas_area_px if canvas_area_px > 0 else 0
            is_full_bleed = (
                area_ratio >= 0.85
                and abs(sx) < canvas_w_px * 0.05
                and abs(sy) < canvas_h_px * 0.05
                and abs(sw - canvas_w_px) < canvas_w_px * 0.05
                and abs(sh - canvas_h_px) < canvas_h_px * 0.05
            )

            is_small = area_ratio < 0.15

            if blob_sha not in sha_assets:
                # Assign asset_id
                if is_full_bleed:
                    asset_counters["bg"] += 1
                    asset_id = f"bg_{asset_counters['bg']:02d}"
                    usage = "full_bleed_background"
                elif is_small:
                    asset_counters["logo"] += 1
                    asset_id = f"logo_{asset_counters['logo']:02d}"
                    usage = "logo"
                else:
                    asset_counters["img"] += 1
                    asset_id = f"img_{asset_counters['img']:02d}"
                    usage = "content"

                # Write file
                fname = f"{asset_id}{ext}"
                dest = media_dir / fname
                dest.write_bytes(blob)

                sha_assets[blob_sha] = {
                    "asset_id": asset_id,
                    "file": fname,
                    "usage": usage,
                    "pages": [],
                    "geometries": [],
                    "blob": blob,
                }

            sha_assets[blob_sha]["pages"].append(slide_idx)
            sha_assets[blob_sha]["geometries"].append({
                "x": sx, "y": sy, "width": sw, "height": sh,
            })

    # Build output list
    assets = []
    for sha, info in sha_assets.items():
        pages = sorted(info["pages"])
        usage = info["usage"]

        # Re-classify if used on >=3 pages and not full_bleed
        if len(pages) >= 3 and usage != "full_bleed_background":
            usage = "decorative"

        # Determine fit strategy
        if usage == "full_bleed_background":
            fit = "cover"
        elif usage == "logo":
            fit = "none"
        else:
            fit = "contain"

        # Average geometry across appearances
        geo = info["geometries"][0]
        avg_x = round(sum(g["x"] for g in info["geometries"]) / len(info["geometries"]), 1)
        avg_y = round(sum(g["y"] for g in info["geometries"]) / len(info["geometries"]), 1)
        avg_w = round(sum(g["width"] for g in info["geometries"]) / len(info["geometries"]), 1)
        avg_h = round(sum(g["height"] for g in info["geometries"]) / len(info["geometries"]), 1)

        confidence = "high"

        assets.append({
            "asset_id": info["asset_id"],
            "file": info["file"],
            "usage": usage,
            "applies_to": {
                "page_types": [],
                "page_indices": pages,
            },
            "source_pages": pages,
            "display_geometry": {
                "x": avg_x,
                "y": avg_y,
                "width": avg_w,
                "height": avg_h,
            },
            "fit": fit,
            "source_canvas": canvas_str,
            "target_canvas": canvas_str,
            "confidence": confidence,
        })

    return assets


# ---------------------------------------------------------------------------
# c.1) extract_native_shape_candidates
# ---------------------------------------------------------------------------

def _shape_fill_hex(shape):
    """Return a direct RGB fill when python-pptx exposes one, otherwise None."""
    try:
        rgb = shape.fill.fore_color.rgb
        return f"#{rgb}" if rgb is not None else None
    except Exception:
        return None


def _shape_line_style(shape):
    """Return factual line styling without inferring theme colors."""
    stroke = None
    width = 0.0
    try:
        rgb = shape.line.color.rgb
        stroke = f"#{rgb}" if rgb is not None else None
    except Exception:
        pass
    try:
        width = emu_to_px(shape.line.width)
    except Exception:
        pass
    return {"stroke": stroke, "stroke_width": width}


def _shape_kind(shape):
    """Return the small, portable vocabulary used by the fidelity builder."""
    if "LINE" in str(getattr(shape, "shape_type", "")).upper():
        return "line"
    try:
        name = str(shape.auto_shape_type).lower()
        if "round" in name and "rect" in name:
            return "roundRect"
        if "rect" in name:
            return "rect"
        if "ellipse" in name or "oval" in name:
            return "ellipse"
        if "line" in name:
            return "line"
    except Exception:
        pass
    return "unsupported"


def _walk_native_shapes(shapes, page_index, prefix="", source_layer="slide", source_pages=None):
    """Collect native auto-shapes without treating them as visual conclusions."""
    candidates = []
    for position, shape in enumerate(shapes, start=1):
        shape_path = f"{prefix}.{position}" if prefix else str(position)
        shape_type = str(getattr(shape, "shape_type", "")).upper()
        if "GROUP" in shape_type and hasattr(shape, "shapes"):
            try:
                candidates.extend(_walk_native_shapes(
                    shape.shapes, page_index, shape_path, source_layer=source_layer, source_pages=source_pages,
                ))
                continue
            except Exception:
                pass
        # Pictures, placeholders, charts and text boxes have separate evidence
        # paths. Treating them as rectangles created false fidelity candidates.
        if not any(name in shape_type for name in ("AUTO_SHAPE", "LINE", "FREEFORM")):
            continue
        kind = _shape_kind(shape)
        if kind == "unsupported":
            continue
        width = emu_to_px(getattr(shape, "width", 0))
        height = emu_to_px(getattr(shape, "height", 0))
        if kind == "line":
            if width <= 0 and height <= 0:
                continue
        elif width <= 0 or height <= 0:
            continue
        style = {"fill": _shape_fill_hex(shape), **_shape_line_style(shape)}
        contains_text = False
        text_length = 0
        try:
            text = str(shape.text or "").strip() if shape.has_text_frame else ""
            contains_text = bool(text)
            text_length = len(text)
        except Exception:
            pass
        layer_prefix = f"slide_{page_index:02d}" if source_layer == "slide" else source_layer
        candidates.append({
            "candidate_id": f"{layer_prefix}.shape_{shape_path.replace('.', '_')}",
            "source_layer": source_layer,
            "source_page": page_index,
            "source_pages": list(source_pages or ([page_index] if page_index else [])),
            "kind": "native_shape",
            "shape": {"preset": kind},
            "geometry": {
                "x": emu_to_px(getattr(shape, "left", 0)),
                "y": emu_to_px(getattr(shape, "top", 0)),
                "width": width,
                "height": height,
            },
            "style": style,
            "contains_text": contains_text,
            "text_length": text_length,
            "shape_name": str(getattr(shape, "name", "")),
            "confidence": "high",
            "requires_visual_confirmation": True,
        })
    return candidates


def _native_shape_repeat_groups(candidates, canvas_width, canvas_height):
    """Group exact/near-exact recurrent candidates as factual evidence."""
    grouped = defaultdict(list)
    for item in candidates:
        geometry = item["geometry"]
        style = item.get("style", {})
        signature = (
            item.get("shape", {}).get("preset"),
            round(geometry["x"] / max(canvas_width, 1), 2),
            round(geometry["y"] / max(canvas_height, 1), 2),
            round(geometry["width"] / max(canvas_width, 1), 2),
            round(geometry["height"] / max(canvas_height, 1), 2),
            style.get("fill"), style.get("stroke"), round(float(style.get("stroke_width") or 0), 1),
            bool(item.get("contains_text")),
        )
        grouped[signature].append(item)
    groups = []
    for index, occurrences in enumerate((items for items in grouped.values() if len(items) >= 2), start=1):
        group_id = f"native_repeat_{index:02d}"
        pages = sorted({page for item in occurrences for page in item.get("source_pages", [])})
        ids = [item["candidate_id"] for item in occurrences]
        for item in occurrences:
            item.update({"repeat_group_id": group_id, "occurrence_count": len(occurrences), "occurrence_pages": pages})
        groups.append({"group_id": group_id, "candidate_ids": ids, "occurrence_count": len(occurrences), "source_pages": pages})
    return groups


def extract_native_shape_candidates(prs):
    """Extract factual native-shape candidates from every slide.

    This intentionally does not decide whether a shape is decoration, content,
    or a reusable component. That decision remains with the visual Template
    Worker in the fidelity-template build step.
    """
    candidates = []
    for page_index, slide in enumerate(prs.slides, start=1):
        candidates.extend(_walk_native_shapes(slide.shapes, page_index))
    for master_index, master in enumerate(prs.slide_masters, start=1):
        master_pages = [i for i, slide in enumerate(prs.slides, start=1) if slide.slide_layout.slide_master == master]
        candidates.extend(_walk_native_shapes(
            master.shapes, 0, source_layer=f"master_{master_index:02d}", source_pages=master_pages,
        ))
        for layout_index, layout in enumerate(master.slide_layouts, start=1):
            layout_pages = [i for i, slide in enumerate(prs.slides, start=1) if slide.slide_layout == layout]
            if layout_pages:
                candidates.extend(_walk_native_shapes(
                    layout.shapes, 0, source_layer=f"layout_{master_index:02d}_{layout_index:02d}", source_pages=layout_pages,
                ))
    canvas = {"width": emu_to_px(prs.slide_width), "height": emu_to_px(prs.slide_height)}
    return candidates, _native_shape_repeat_groups(candidates, canvas["width"], canvas["height"]), canvas


# ---------------------------------------------------------------------------
# d) extract_font_sizes_from_svg
# ---------------------------------------------------------------------------

def find_svg_flat_dir(project_dir):
    """Find svg-flat/ directory under project_dir."""
    candidates = [
        Path(project_dir) / "_internal" / "00_project" / "template_visuals" / "svg-flat",
        Path(project_dir) / "_internal" / "00_project" / "svg-flat",
        Path(project_dir) / "template_visuals" / "svg-flat",
    ]
    for c in candidates:
        if c.is_dir():
            return c
    return None


def _parse_svg_text_sizes(svg_path):
    """Parse font sizes from SVG text elements. Returns list of (size_px, font_family, font_weight)."""
    results = []
    try:
        tree = ET.parse(svg_path)
        root = tree.getroot()
    except Exception:
        return results

    # Namespace handling for SVG
    ns = {"svg": "http://www.w3.org/2000/svg"}

    for text_el in root.iter("{http://www.w3.org/2000/svg}text"):
        style = text_el.get("style", "") or ""
        font_size_str = text_el.get("font-size", "")

        # Parse from style attribute
        if not font_size_str and "font-size:" in style:
            m = re.search(r'font-size:\s*([\d.]+)(px|pt|)', style)
            if m:
                font_size_str = m.group(1) + (m.group(2) or "px")

        if not font_size_str:
            continue

        # Convert to px
        size_px = None
        m = re.match(r'([\d.]+)\s*px', font_size_str)
        if m:
            size_px = float(m.group(1))
        m = re.match(r'([\d.]+)\s*pt', font_size_str)
        if m:
            size_px = float(m.group(1)) * 1.333  # pt -> px
        m = re.match(r'^([\d.]+)$', font_size_str)
        if m:
            size_px = float(m.group(1))

        if size_px is None or size_px <= 0 or size_px > 200:
            continue

        # Font family
        font_family = text_el.get("font-family", "")
        if not font_family and "font-family:" in style:
            m = re.search(r'font-family:\s*([^;]+)', style)
            if m:
                font_family = m.group(1).strip().strip("'\"")

        # Font weight
        font_weight = text_el.get("font-weight", "")
        if not font_weight and "font-weight:" in style:
            m = re.search(r'font-weight:\s*([^;]+)', style)
            if m:
                font_weight = m.group(1).strip()
        if not font_weight:
            font_weight = "normal"

        results.append((size_px, font_family or "", font_weight))

    return results


def _cluster_font_sizes(sizes_by_page):
    """
    Cluster font sizes across pages into roles.
    Returns list of {role, font_family, font_size_px, font_weight, source, confidence}.
    """
    if not sizes_by_page:
        return []

    # Collect all unique sizes across pages
    all_sizes = []
    for page_sizes in sizes_by_page.values():
        all_sizes.extend(page_sizes)

    if not all_sizes:
        return []

    # Get unique (size_px, font_family, font_weight) tuples
    unique = list(set(all_sizes))
    unique.sort(key=lambda x: x[0], reverse=True)

    # Assign roles by size rank
    roles_list = ["cover_title", "heading", "subheading", "body", "caption"]
    role_results = []

    for i, (size_px, font_family, font_weight) in enumerate(unique):
        if i >= len(roles_list):
            role = f"level_{i + 1}"
        else:
            role = roles_list[i]

        # Confidence: if this size appears on >=3 pages, high
        count = sum(1 for s in all_sizes if abs(s[0] - size_px) < 1)
        total_pages = len(sizes_by_page)
        confidence = "high" if count >= min(3, total_pages) else "medium"

        entry = {
            "role": role,
            "font_family": font_family,
            "font_size_px": int(size_px),
            "font_weight": font_weight,
            "source": "svg_measured",
            "confidence": confidence,
        }
        role_results.append(entry)

    return role_results


def extract_font_sizes_from_svg(project_dir):
    """
    Look for svg-flat/ slides, parse SVG text elements, extract font-size values.
    Cluster by size (largest=cover_title, then heading, body, caption).
    Returns type_hierarchy supplement list.
    Returns empty list if svg-flat is not available.
    """
    svg_dir = find_svg_flat_dir(project_dir)
    if svg_dir is None:
        return []

    svg_files = sorted(svg_dir.glob("slide_*.svg"))
    # Also try pattern without slide_ prefix
    if not svg_files:
        svg_files = sorted(svg_dir.glob("*.svg"))

    if not svg_files:
        return []

    sizes_by_page = {}
    for svg_path in svg_files:
        page_key = svg_path.stem
        sizes = _parse_svg_text_sizes(svg_path)
        if sizes:
            sizes_by_page[page_key] = sizes

    return _cluster_font_sizes(sizes_by_page)


# ---------------------------------------------------------------------------
# e) extract_decoration_patterns_from_svg
# ---------------------------------------------------------------------------

_DECORATION_SHAPE_TYPES = {
    "{http://www.w3.org/2000/svg}rect",
    "{http://www.w3.org/2000/svg}line",
    "{http://www.w3.org/2000/svg}circle",
    "{http://www.w3.org/2000/svg}ellipse",
}


def _is_white_fill(el):
    """Check if an SVG element has white fill."""
    fill = el.get("fill", "").lower().strip()
    if fill in ("white", "#ffffff", "#fff", "rgb(255,255,255)", "rgba(255,255,255,1)"):
        return True
    style = el.get("style", "")
    m = re.search(r'fill:\s*(#[0-9a-fA-F]+|white|rgb\(255,\s*255,\s*255\))', style)
    if m:
        val = m.group(1).lower()
        if val in ("white", "#ffffff", "#fff"):
            return True
    return False


def _extract_decoration_elements(svg_path):
    """Extract non-text rect/line elements from a single SVG page."""
    elements = []
    try:
        tree = ET.parse(svg_path)
        root = tree.getroot()
    except Exception:
        return elements

    ns = "http://www.w3.org/2000/svg"
    viewbox = root.get("viewBox", "")
    vb_parts = viewbox.split()
    canvas_w = float(vb_parts[2]) if len(vb_parts) >= 3 else 1280
    canvas_h = float(vb_parts[3]) if len(vb_parts) >= 4 else 720

    for el in root.iter():
        tag = el.tag
        if tag not in _DECORATION_SHAPE_TYPES:
            continue

        # Skip white fills
        if _is_white_fill(el):
            continue

        # Skip elements with no fill and no stroke
        fill = el.get("fill", el.get("style", "") or "").lower()
        stroke = el.get("stroke", el.get("style", "") or "").lower()
        if "none" in fill and "none" in stroke:
            continue

        # Determine geometry for filtering
        if tag == "{http://www.w3.org/2000/svg}rect":
            x = float(el.get("x", 0))
            y = float(el.get("y", 0))
            w = float(el.get("width", canvas_w))
            h = float(el.get("height", canvas_h))
        elif tag == "{http://www.w3.org/2000/svg}line":
            x = min(float(el.get("x1", 0)), float(el.get("x2", 0)))
            y = min(float(el.get("y1", 0)), float(el.get("y2", 0)))
            w = abs(float(el.get("x2", 0)) - float(el.get("x1", 0))) or 1
            h = abs(float(el.get("y2", 0)) - float(el.get("y1", 0))) or 1
        else:
            cx = float(el.get("cx", canvas_w / 2))
            cy = float(el.get("cy", canvas_h / 2))
            r = float(el.get("r", 0))
            x = cx - r
            y = cy - r
            w = r * 2
            h = r * 2

        # Filter out full-canvas elements (>80%)
        area_ratio = (w * h) / (canvas_w * canvas_h) if canvas_w * canvas_h > 0 else 0
        if area_ratio > 0.8:
            continue

        # Determine position zone
        zone_x = "left" if x < canvas_w / 3 else "center" if x < canvas_w * 2 / 3 else "right"
        zone_y = "top" if y < canvas_h / 3 else "middle" if y < canvas_h * 2 / 3 else "bottom"

        # Color key
        color = el.get("fill", "")
        if not color or color == "none":
            color = el.get("stroke", "")
        color = color.strip().lower()

        elem_type = tag.split("}")[-1]  # strip namespace

        elements.append({
            "type": elem_type,
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "fill": el.get("fill", "none"),
            "stroke": el.get("stroke", "none"),
            "stroke_width": el.get("stroke-width", "1"),
            "rx": el.get("rx", "0"),
            "color_key": color,
            "zone_x": zone_x,
            "zone_y": zone_y,
            "position_zone": f"{zone_x}-{zone_y}",
        })

    return elements


def extract_decoration_patterns_from_svg(project_dir):
    """
    Parse svg-flat SVGs. Find non-text rect/line elements.
    Filter out full-canvas elements (>80%), white fills.
    Cluster by (color, type, position). Patterns on >=2 pages = output.

    Returns decoration_patterns[] list.
    Returns empty list if svg-flat not available.
    """
    svg_dir = find_svg_flat_dir(project_dir)
    if svg_dir is None:
        return []

    svg_files = sorted(svg_dir.glob("slide_*.svg"))
    if not svg_files:
        svg_files = sorted(svg_dir.glob("*.svg"))
    if not svg_files:
        return []

    # Extract elements per page
    page_elements = {}
    for svg_path in svg_files:
        page_key = svg_path.stem
        elements = _extract_decoration_elements(svg_path)
        if elements:
            page_elements[page_key] = elements

    if not page_elements:
        return []

    # Cluster by color_key + type + position_zone
    cluster_map = defaultdict(list)  # (color, type, zone) -> [(page_key, elem)]
    for page_key, elements in page_elements.items():
        # Dedup within a page: keep one element per (type, color, zone) combo
        seen_in_page = set()
        for elem in elements:
            cluster_key = (elem["color_key"], elem["type"], elem["position_zone"])
            if cluster_key in seen_in_page:
                continue
            seen_in_page.add(cluster_key)
            cluster_map[cluster_key].append((page_key, elem))

    # Filter: patterns on >=2 pages
    decoration_patterns = []
    pattern_counter = 0

    for cluster_key, instances in cluster_map.items():
        unique_pages = set(pk for pk, _ in instances)
        if len(unique_pages) < 2:
            continue

        pattern_counter += 1
        color, elem_type, zone = cluster_key

        # Build elements list (dedup by type within the pattern)
        pattern_elements = []
        seen_shapes = set()
        for page_key, elem in instances:
            shape_key = (
                elem["type"], elem["fill"], elem["stroke"],
                round(elem["x"], -1), round(elem["y"], -1),
            )
            if shape_key in seen_shapes:
                continue
            seen_shapes.add(shape_key)

            pe = {
                "type": elem["type"],
                "x": round(elem["x"], 1),
                "y": round(elem["y"], 1),
                "width": round(elem["width"], 1),
                "height": round(elem["height"], 1),
            }
            if elem["fill"] and elem["fill"] != "none":
                pe["fill"] = elem["fill"]
            if elem["stroke"] and elem["stroke"] != "none":
                pe["stroke"] = elem["stroke"]
                pe["stroke-width"] = elem["stroke_width"]
            if elem["type"] == "rect" and elem["rx"] and elem["rx"] != "0":
                pe["rx"] = float(elem["rx"])
            pattern_elements.append(pe)

        if not pattern_elements:
            continue

        # Determine canvas_change_strategy
        has_full_width = any(
            e["type"] == "rect" and e["width"] >= 1000
            for e in pattern_elements
        )
        strategy = "stretch_width" if has_full_width else "scale_y"

        min_x = min(e["x"] for e in pattern_elements)
        if min_x <= 0 and has_full_width:
            strategy = "stretch_width"
        elif min_x > 0:
            strategy = "scale_y"

        # Determine name and page_types from zone
        color_name = color.lstrip("#")[:6] if color.startswith("#") else color
        type_label = elem_type.capitalize()
        zone_label = zone.replace("-", "_")

        decoration_patterns.append({
            "pattern_id": f"pattern_{pattern_counter:02d}",
            "name": f"{type_label} at {zone_label}",
            "applies_to": {
                "page_types": [],
            },
            "elements": pattern_elements,
            "parameters": {
                "suggested_scale": 1.0,
                "canvas_change_strategy": strategy,
            },
            "confidence": "high" if len(unique_pages) >= 3 else "medium",
        })

    return decoration_patterns


# ---------------------------------------------------------------------------
# f) merge_into_profile
# ---------------------------------------------------------------------------

def merge_into_profile(existing, colors, fonts, assets, font_sizes, patterns, native_shapes, native_shape_groups, canvas):
    """
    Preserve ALL existing fields in template_profile.json.
    Preserve visual conclusions. Store PPTX/XML output as structural candidates
    for the visual Template Worker to confirm or reject.
    """
    profile = dict(existing) if existing else {}

    # Ensure source_files exists
    if "source_files" not in profile:
        profile["source_files"] = []
    if "pages_reviewed" not in profile:
        profile["pages_reviewed"] = []
    if "method" not in profile:
        profile["method"] = "visual_only"
    if "limitations" not in profile:
        profile["limitations"] = []
    if "generated_at" not in profile:
        profile["generated_at"] = ""

    # Ensure design_direction exists
    if "design_direction" not in profile or not isinstance(profile["design_direction"], dict):
        profile["design_direction"] = {}

    dd = profile["design_direction"]

    # Ensure default fields in design_direction
    _ensure_dd_field(dd, "overall_character", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "title_entry", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "grid_and_alignment", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "spacing_and_density", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "image_language", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "chart_language", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "component_language", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "deck_rhythm", {"value": "", "evidence_pages": [], "confidence": "low"})
    _ensure_dd_field(dd, "reusable_motifs", [])
    _ensure_dd_field(dd, "page_exceptions", [])

    _ensure_dd_field(dd, "color_roles", [])
    _ensure_dd_field(dd, "type_hierarchy", [])
    # Migrate profiles made by the earlier extractor: XML alone is not a
    # visual conclusion and must not remain in the worker-facing direction.
    dd["color_roles"] = [item for item in dd["color_roles"] if item.get("source") != "xml_parsed"]
    dd["type_hierarchy"] = [item for item in dd["type_hierarchy"] if item.get("source") != "xml_parsed"]
    has_visual_conclusions = (
        bool(str(dd.get("overall_character", {}).get("value", "")).strip())
        or bool(dd["color_roles"])
        or bool(dd["type_hierarchy"])
    )
    if has_visual_conclusions:
        profile.setdefault("reusable_assets", [])
        profile.setdefault("decoration_patterns", [])
    else:
        # Remove unsafe output from the earlier binding extractor. A visual
        # Worker may repopulate these fields after checking every page.
        profile["pages_reviewed"] = []
        profile["reusable_assets"] = []
        profile["decoration_patterns"] = []

    # XML/PPTX data is evidence, not a visual conclusion. Full-slide pictures
    # can contain placeholder copy, so assets must be accepted by vision first.
    profile["structural_extraction"] = {
        "colors": colors,
        "fonts": fonts,
        "assets": assets,
        "font_sizes": font_sizes,
        "decoration_patterns": patterns,
        "native_shapes": native_shapes,
        "native_shape_groups": native_shape_groups,
        "canvas": canvas,
        "requires_visual_confirmation": True,
    }

    if not has_visual_conclusions:
        profile["usage_policy"] = {
            "mode": "extraction_audit_only",
            "note": "XML 与 PPTX 图片仅为结构候选；完整 profile 不进入 SVG task。",
        }

    # Update generated_at
    profile["generated_at"] = datetime.now(timezone.utc).isoformat()

    return profile


def _ensure_dd_field(dd, key, default):
    """Ensure design_direction dict has a field with default if missing."""
    if key not in dd or dd[key] is None:
        dd[key] = default


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Extract precise design tokens, bitmap assets, and decoration patterns from PPTX into template_profile.json"
    )
    parser.add_argument("source", help="Path to the PPTX template file")
    parser.add_argument("--project", required=True, help="Project root directory")
    args = parser.parse_args()

    source = Path(args.source)
    project_dir = Path(args.project).resolve()

    if not source.exists():
        print(f"ERROR: PPTX file not found: {source}", file=sys.stderr)
        sys.exit(1)

    if source.suffix.lower() not in (".pptx", ".pptm"):
        print(f"ERROR: Not a PPTX file: {source}", file=sys.stderr)
        sys.exit(1)

    internal_dir = project_dir / "_internal" / "00_project"
    media_dir = internal_dir / "template_media"
    profile_path = internal_dir / "template_profile.json"

    # 1. Read existing template_profile.json
    existing = {}
    if profile_path.exists():
        try:
            existing = json.loads(profile_path.read_text(encoding="utf-8"))
            print(f"Read existing template_profile.json: {profile_path}")
        except (json.JSONDecodeError, Exception) as e:
            print(f"Warning: Could not parse existing template_profile.json ({e}), starting fresh")

    # 2. Open PPTX and extract
    print(f"Analyzing: {source.name}")
    prs = Presentation(str(source))

    print("  Extracting theme colors from all masters...")
    colors = extract_theme_colors_all_masters(prs, pptx_path=source)
    print(f"    Found {len(colors)} color roles")

    print("  Extracting theme fonts from all masters...")
    fonts = extract_theme_fonts_all_masters(prs, pptx_path=source)
    print(f"    Found {len(fonts)} font roles")

    print("  Extracting bitmap assets from all slides...")
    assets = extract_assets_from_slides(prs, project_dir)
    print(f"    Found {len(assets)} reusable assets")

    print("  Extracting native shape candidates from all slides...")
    native_shapes, native_shape_groups, canvas = extract_native_shape_candidates(prs)
    print(f"    Found {len(native_shapes)} native shape candidates")

    # 3. Extract from svg-flat if available
    print("  Checking for svg-flat directory...")
    font_sizes = extract_font_sizes_from_svg(project_dir)
    if font_sizes:
        print(f"    Measured {len(font_sizes)} font sizes from SVG")
    else:
        print("    svg-flat not available, skipping font size measurement")

    patterns = extract_decoration_patterns_from_svg(project_dir)
    if patterns:
        print(f"    Found {len(patterns)} decoration patterns from SVG")
    else:
        print("    svg-flat not available, skipping decoration pattern extraction")

    # 4. Merge into profile
    print("  Merging into template_profile.json...")
    enhanced = merge_into_profile(existing, colors, fonts, assets, font_sizes, patterns, native_shapes, native_shape_groups, canvas)

    # Source provenance is safe to populate. Page review ownership stays with
    # the visual Template Worker; structural extraction must not claim it.
    if not enhanced.get("source_files"):
        enhanced["source_files"] = [str(source.resolve())]

    # 5. Write output
    profile_path.parent.mkdir(parents=True, exist_ok=True)
    profile_path.write_text(json.dumps(enhanced, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"\nTemplate profile enhanced: {profile_path}")
    print(f"  Color candidates: {len(colors)} (XML evidence; visual confirmation required)")
    print(f"  Font candidates: {len(fonts)} (XML evidence; visual confirmation required)")
    print(f"  Bitmap asset candidates: {len(assets)}")
    print(f"  Font size entries: {len(font_sizes)}")
    print(f"  Decoration patterns: {len(patterns)}")
    print(f"  Native shape candidates: {len(native_shapes)}")
    print(f"  Repeated native shape groups: {len(native_shape_groups)}")
    print(f"  Usage policy mode: reference_only")
    print(f"  Template media: {media_dir}")


if __name__ == "__main__":
    main()
